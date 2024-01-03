from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageFilter

def find_slide_aspect_ratio(image_paths):
    aspect_ratio = False
    for image_path in image_paths:
        image = Image.open(image_path)
        width_height = image.width/image.height
        if not aspect_ratio:
            aspect_ratio = width_height
        else:
            if abs(1-width_height/aspect_ratio) > 0.01:
                raise Exception('Incosistent aspect ratio!')
    return aspect_ratio

def apply_gaussian_blur(image_path, blur_radius=5):
    new_path = image_path.split('.')[0]
    new_path = f'{image_path}_blurred.jpg'

    image = Image.open(image_path)
    img_blurred = image.filter(ImageFilter.GaussianBlur(radius=blur_radius))
    img_blurred.save(new_path)

    return new_path

def create_powerpoint(image_paths,output_path,aspect_ratio=1.7777):
    # Load the PowerPoint presentation
    presentation = Presentation()
    presentation.slide_height=Inches(10)
    presentation.slide_width=Inches(10*aspect_ratio)

    # Iterate through each slide in the presentation
    for image_file in image_paths:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Use a blank slide layout (index 5)

        # Add image to the slide
        left = top = Inches(0)
        pic = slide.shapes.add_picture(image_file, left, top, width=Inches(10*aspect_ratio))

    # Save the modified presentation
    presentation.save(output_path)

def make_powerpoint_from_image_files(image_files):
    image_files.sort()

    aspect_ratio = find_slide_aspect_ratio(image_files)

    blurred_files = []
    for i in image_files:
        blurred_files.append(apply_gaussian_blur(i))

    output_path = '/blurred_powerpoint.pptx'
    create_powerpoint(blurred_files,output_path,aspect_ratio=aspect_ratio)

    return [output_path]

def js_entry_point(image_files):
    return make_powerpoint_from_image_files(image_files.to_py())


if __name__ == "__main__":
    import glob
    input_folder_path = "test_powerpoint"

    image_files = glob.glob(f'{input_folder_path}/*')
    make_powerpoint_from_image_files(image_files)
